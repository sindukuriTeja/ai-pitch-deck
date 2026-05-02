import html.parser
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.pptx_engine.themes import Theme, rgb

class PPTXHTMLParser(html.parser.HTMLParser):
    def __init__(self, slide, theme: Theme, is_title_slide: bool = False, image_path: str = None):
        super().__init__()
        self.slide = slide
        self.theme = theme
        self.is_title_slide = is_title_slide
        self.image_path = image_path
        self.current_tag = None
        
        # Slide Dimensions
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        
        # Add background
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), self.slide_width, self.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(theme.colors.bg_dark)
        shape.line.fill.background()

        # Handle Image Placement Correctly
        margin_x = Inches(0.8)
        content_width = self.slide_width - (2 * margin_x)
        
        if self.image_path and os.path.exists(self.image_path):
            # Two-column layout: Text on left (60%), Image on right (35%)
            text_width = content_width * 0.55
            image_width = content_width * 0.40
            
            # Add Image with border/frame styling
            img_left = margin_x + text_width + Inches(0.5)
            img_top = Inches(1.5)
            # Try to add the picture
            try:
                self.slide.shapes.add_picture(self.image_path, img_left, img_top, width=image_width)
            except Exception:
                # If image loading fails, revert to full width
                text_width = content_width
            
            self.text_frame_width = text_width
            self.text_left = margin_x
        else:
            self.text_frame_width = content_width
            self.text_left = margin_x

        # Vertical Positioning
        if self.is_title_slide:
            top = Inches(2.2)
            height = Inches(3.5)
            # If title slide has image, keep text left. If not, center it.
            self.text_alignment = PP_ALIGN.LEFT if self.image_path else PP_ALIGN.CENTER
        else:
            top = Inches(0.8)
            height = Inches(6.0)
            self.text_alignment = PP_ALIGN.LEFT
            
        self.txBox = self.slide.shapes.add_textbox(
            self.text_left, top, self.text_frame_width, height
        )
        self.tf = self.txBox.text_frame
        self.tf.word_wrap = True
        
        if self.is_title_slide:
            self.tf.vertical_anchor = 1 # Top

    def handle_starttag(self, tag, attrs):
        self.current_tag = tag

    def handle_endtag(self, tag):
        self.current_tag = None

    def handle_data(self, data):
        text = data.strip()
        if not text:
            return

        if self.current_tag == 'h1':
            p = self._get_new_paragraph()
            p.text = text
            p.font.name = self.theme.fonts.heading
            p.font.size = Pt(60 if self.is_title_slide else 42)
            p.font.color.rgb = rgb(self.theme.colors.text_light)
            p.font.bold = True
            p.alignment = self.text_alignment
            p.space_after = Pt(24)
            
        elif self.current_tag == 'h2':
            p = self._get_new_paragraph()
            p.text = text
            p.font.name = self.theme.fonts.heading
            p.font.size = Pt(32 if self.is_title_slide else 26)
            p.font.color.rgb = rgb(self.theme.colors.accent)
            p.font.bold = True
            p.alignment = self.text_alignment
            p.space_after = Pt(18)
            
        elif self.current_tag == 'h3':
            p = self._get_new_paragraph()
            p.text = text
            p.font.name = self.theme.fonts.heading
            p.font.size = Pt(24)
            p.font.color.rgb = rgb(self.theme.colors.text_light)
            p.font.bold = True
            p.space_before = Pt(12)
            p.space_after = Pt(6)
            
        elif self.current_tag == 'p':
            p = self._get_new_paragraph()
            p.text = text
            p.font.name = self.theme.fonts.body
            p.font.size = Pt(18)
            p.font.color.rgb = rgb(self.theme.colors.text_light)
            p.space_before = Pt(12)
            p.alignment = self.text_alignment
            
        elif self.current_tag == 'li':
            p = self._get_new_paragraph()
            p.text = f"• {text}"
            p.font.name = self.theme.fonts.body
            p.font.size = Pt(18)
            p.font.color.rgb = rgb(self.theme.colors.text_light)
            p.level = 0
            p.space_before = Pt(6)

    def _get_new_paragraph(self):
        if len(self.tf.paragraphs) == 1 and not self.tf.paragraphs[0].text:
            return self.tf.paragraphs[0]
        return self.tf.add_paragraph()

def parse_html_to_slide(slide, html_str: str, theme: Theme, is_title_slide: bool = False, image_path: str = None):
    parser = PPTXHTMLParser(slide, theme, is_title_slide, image_path)
    parser.feed(html_str)
