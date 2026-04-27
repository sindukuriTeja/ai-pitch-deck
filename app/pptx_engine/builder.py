import os
from pptx import Presentation
from pptx.util import Inches
from app.pptx_engine.themes import get_theme, Theme
from app.pptx_engine.html_parser import parse_html_to_slide
from app.config import OUTPUT_DIR


def build_presentation(deck_plan: dict, theme_id: str, task_id: str) -> str:
    theme = get_theme(theme_id)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    slides_data = deck_plan.get("slides", [])

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        html_content = slide_data.get("html", "<h1>Missing Content</h1>")
        
        try:
            parse_html_to_slide(slide, html_content, theme)
        except Exception as e:
            # Fallback if parser fails
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(5))
            txBox.text_frame.text = "Error parsing HTML: " + str(e)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, f"{task_id}.pptx")
    prs.save(output_path)
    return output_path
