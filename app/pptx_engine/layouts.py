from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from app.pptx_engine.themes import Theme, rgb


def _add_shape(slide, left, top, width, height, fill_color: str, rounded: bool = False):
    shape_type = 5 if rounded else 1
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.fill.background()
    if rounded and shape.adjustments:
        shape.adjustments[0] = 0.05
    return shape


def _add_textbox(slide, left, top, width, height, text, font_name, font_size, color, bold=False, alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = rgb(color)
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_bullets(slide, left, top, width, height, bullets, font_name, font_size, color):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(12) # Increased spacing for modern look
    return txBox


# ---- LAYOUT BUILDERS ----

def build_title_slide(slide, content: dict, theme: Theme):
    """Full background color with centered title and subtitle."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_dark)
    # Elegant top accent bar
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), theme.colors.accent)

    _add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.33), Inches(1.5),
                 content.get("title", ""), theme.fonts.heading, 52,
                 theme.colors.text_light, bold=True, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.BOTTOM)

    if content.get("subtitle"):
        _add_textbox(slide, Inches(2), Inches(4.0), Inches(9.33), Inches(1),
                     content["subtitle"], theme.fonts.body, theme.fonts.subtitle_size + 4,
                     theme.colors.accent, alignment=PP_ALIGN.CENTER)

    if content.get("cta"):
        # CTA Pill
        _add_shape(slide, Inches(4.66), Inches(5.5), Inches(4), Inches(0.8), theme.colors.accent, rounded=True)
        _add_textbox(slide, Inches(4.66), Inches(5.55), Inches(4), Inches(0.7),
                     content["cta"], theme.fonts.body, 16,
                     theme.colors.text_dark, bold=True, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

def build_section_header(slide, content: dict, theme: Theme):
    """Left accent bar with large heading."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_light)
    
    # Modern rounded block on left
    _add_shape(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(4.5), theme.colors.primary, rounded=True)

    _add_textbox(slide, Inches(1), Inches(2.5), Inches(4), Inches(2.5),
                 content.get("title", ""), theme.fonts.heading, 42,
                 theme.colors.text_light, bold=True, vertical_anchor=MSO_ANCHOR.MIDDLE)

    if content.get("body"):
        _add_textbox(slide, Inches(6.2), Inches(2.5), Inches(6), Inches(3),
                     content["body"], theme.fonts.body, theme.fonts.body_size + 2,
                     theme.colors.text_dark, vertical_anchor=MSO_ANCHOR.MIDDLE)

    if content.get("subtitle"):
        _add_textbox(slide, Inches(6.2), Inches(1.5), Inches(6), Inches(0.8),
                     content["subtitle"].upper(), theme.fonts.body, theme.fonts.subtitle_size - 2,
                     theme.colors.accent, bold=True)

def build_content_image(slide, content: dict, theme: Theme):
    """Text on left, colored placeholder on right."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_light)

    # Title
    _add_textbox(slide, Inches(1), Inches(0.8), Inches(6), Inches(1),
                 content.get("title", ""), theme.fonts.heading, theme.fonts.heading_size,
                 theme.colors.text_dark, bold=True)

    # Accent underline
    _add_shape(slide, Inches(1), Inches(1.8), Inches(1.5), Inches(0.06), theme.colors.accent)

    # Body text
    if content.get("body"):
        _add_textbox(slide, Inches(1), Inches(2.2), Inches(5.5), Inches(2),
                     content["body"], theme.fonts.body, theme.fonts.body_size,
                     theme.colors.text_dark)

    if content.get("bullets"):
        top = Inches(4.5) if content.get("body") else Inches(2.2)
        _add_bullets(slide, Inches(1), top, Inches(5.5), Inches(2.5),
                     content["bullets"], theme.fonts.body, theme.fonts.body_size,
                     theme.colors.text_dark)

    # Image placeholder (rounded rectangle)
    _add_shape(slide, Inches(7.5), Inches(1), Inches(5), Inches(5.5), theme.colors.secondary, rounded=True)
    _add_textbox(slide, Inches(7.5), Inches(3.2), Inches(5), Inches(1),
                 content.get("visual_suggestion", "Visual Placeholder"), theme.fonts.body, 14,
                 theme.colors.text_light, alignment=PP_ALIGN.CENTER)

def build_stats_numbers(slide, content: dict, theme: Theme):
    """3-4 big stat numbers with labels."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_dark)

    _add_textbox(slide, Inches(1), Inches(0.8), Inches(11.33), Inches(1),
                 content.get("title", ""), theme.fonts.heading, theme.fonts.heading_size,
                 theme.colors.text_light, bold=True, alignment=PP_ALIGN.CENTER)

    # Accent line
    _add_shape(slide, Inches(5.66), Inches(1.8), Inches(2), Inches(0.06), theme.colors.accent)

    bullets = content.get("bullets", [])
    if not bullets and content.get("body"):
        bullets = [s.strip() for s in content["body"].split(".") if s.strip()][:4]

    num_stats = min(len(bullets), 4) if bullets else 3
    if num_stats == 0:
        num_stats = 3
        bullets = ["Key Metric 1", "Key Metric 2", "Key Metric 3"]

    box_width = 3.0 if num_stats <= 3 else 2.5
    gap = 0.5
    total_width = num_stats * box_width + (num_stats - 1) * gap
    start_x = (13.33 - total_width) / 2

    for i in range(num_stats):
        x = start_x + i * (box_width + gap)
        # Stat card (Rounded)
        _add_shape(slide, Inches(x), Inches(2.8), Inches(box_width), Inches(3.5), theme.colors.secondary, rounded=True)
        # Top Accent strip on card
        # _add_shape(slide, Inches(x), Inches(2.8), Inches(box_width), Inches(0.1), theme.colors.accent)

        # Number
        _add_textbox(slide, Inches(x), Inches(3.4), Inches(box_width), Inches(1.2),
                     f"#{i+1}", theme.fonts.heading, 48,
                     theme.colors.accent, bold=True, alignment=PP_ALIGN.CENTER)

        # Label
        label = bullets[i] if i < len(bullets) else f"Metric {i+1}"
        _add_textbox(slide, Inches(x + 0.2), Inches(4.8), Inches(box_width - 0.4), Inches(1.2),
                     label, theme.fonts.body, 14,
                     theme.colors.text_light, alignment=PP_ALIGN.CENTER)

def build_bullet_points(slide, content: dict, theme: Theme):
    """Clean bullet list with elegant icons/dots."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_light)

    _add_textbox(slide, Inches(1), Inches(0.8), Inches(11.33), Inches(1),
                 content.get("title", ""), theme.fonts.heading, theme.fonts.heading_size,
                 theme.colors.text_dark, bold=True)

    _add_shape(slide, Inches(1), Inches(1.8), Inches(1.5), Inches(0.06), theme.colors.accent)

    if content.get("body"):
        _add_textbox(slide, Inches(1), Inches(2.2), Inches(11.33), Inches(1),
                     content["body"], theme.fonts.body, theme.fonts.body_size + 2,
                     theme.colors.text_dark)

    bullets = content.get("bullets", [])
    if bullets:
        start_y = 3.5 if content.get("body") else 2.5
        for i, bullet in enumerate(bullets[:5]):
            y = start_y + i * 0.8
            # Elegant rounded dot
            _add_shape(slide, Inches(1.2), Inches(y + 0.1), Inches(0.15), Inches(0.15), theme.colors.accent, rounded=True)
            _add_textbox(slide, Inches(1.6), Inches(y), Inches(10), Inches(0.6),
                         bullet, theme.fonts.body, theme.fonts.body_size + 2,
                         theme.colors.text_dark)

def build_two_column(slide, content: dict, theme: Theme):
    """Split layout with two columns."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_light)

    _add_textbox(slide, Inches(1), Inches(0.8), Inches(11.33), Inches(1),
                 content.get("title", ""), theme.fonts.heading, theme.fonts.heading_size,
                 theme.colors.text_dark, bold=True, alignment=PP_ALIGN.CENTER)

    _add_shape(slide, Inches(5.66), Inches(1.8), Inches(2), Inches(0.06), theme.colors.accent)

    # Left column (rounded)
    _add_shape(slide, Inches(1), Inches(2.5), Inches(5.3), Inches(4.2), theme.colors.primary, rounded=True)

    left_text = content.get("body", "")
    _add_textbox(slide, Inches(1.4), Inches(2.9), Inches(4.5), Inches(3.4),
                 left_text, theme.fonts.body, theme.fonts.body_size + 2,
                 theme.colors.text_light)

    # Right column (rounded)
    _add_shape(slide, Inches(7), Inches(2.5), Inches(5.3), Inches(4.2), theme.colors.secondary, rounded=True)

    bullets = content.get("bullets", [])
    if bullets:
        _add_bullets(slide, Inches(7.4), Inches(2.9), Inches(4.5), Inches(3.4),
                     bullets, theme.fonts.body, theme.fonts.body_size,
                     theme.colors.text_light)

def build_quote_insight(slide, content: dict, theme: Theme):
    """Large centered quote/insight text."""
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_dark)

    # Large aesthetic quote marks
    _add_textbox(slide, Inches(1.5), Inches(0.5), Inches(2), Inches(2),
                 '"', theme.fonts.heading, 140,
                 theme.colors.accent, bold=True)

    # Quote text
    quote_text = content.get("body") or content.get("title", "")
    _add_textbox(slide, Inches(2), Inches(2.5), Inches(9.33), Inches(2.5),
                 quote_text, theme.fonts.heading, 36,
                 theme.colors.text_light, bold=True, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    if content.get("subtitle"):
        _add_textbox(slide, Inches(3), Inches(5.5), Inches(7.33), Inches(0.8),
                     f"- {content['subtitle']} -", theme.fonts.body, theme.fonts.subtitle_size + 2,
                     theme.colors.accent, alignment=PP_ALIGN.CENTER)

def build_full_image(slide, content: dict, theme: Theme):
    """Full background with elegant overlay text."""
    # Background
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.secondary)
    
    # Elegant dark semi-transparent overlay (faked with solid dark for now)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), theme.colors.bg_dark)
    
    _add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.33), Inches(1.5),
                 content.get("title", ""), theme.fonts.heading, 52,
                 theme.colors.text_light, bold=True, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.BOTTOM)

    if content.get("body"):
        _add_textbox(slide, Inches(2), Inches(4.2), Inches(9.33), Inches(1.5),
                     content["body"], theme.fonts.body, theme.fonts.body_size + 4,
                     theme.colors.text_light, alignment=PP_ALIGN.CENTER)

    if content.get("cta"):
        # Large CTA pill
        _add_shape(slide, Inches(4.66), Inches(6.0), Inches(4), Inches(0.8), theme.colors.accent, rounded=True)
        _add_textbox(slide, Inches(4.66), Inches(6.05), Inches(4), Inches(0.7),
                     content["cta"], theme.fonts.body, 16,
                     theme.colors.text_dark, bold=True, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)


LAYOUT_BUILDERS = {
    "title_slide": build_title_slide,
    "section_header": build_section_header,
    "content_image": build_content_image,
    "stats_numbers": build_stats_numbers,
    "bullet_points": build_bullet_points,
    "two_column": build_two_column,
    "quote_insight": build_quote_insight,
    "full_image": build_full_image,
}
